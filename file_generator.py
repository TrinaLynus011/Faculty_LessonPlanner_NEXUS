"""
File Generator - Creates actual PPT, PDF, Audio, and Video files
Similar to GPT and NotebookLM output
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER
from gtts import gTTS
import json
from PIL import Image, ImageDraw, ImageFont
import subprocess


class FileGenerator:
    """Generates actual files from content data"""
    
    def __init__(self, output_dir="generated_files"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
        self.use_topic_folders = True  # Organize by topic
    
    def generate_all(self, content_data, subject_name):
        """Generate all file types for the content"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_subject = self._sanitize_filename(subject_name)
        
        generated_files = {
            "subject": subject_name,
            "timestamp": timestamp,
            "files": []
        }
        
        # Generate files for each topic
        for idx, content in enumerate(content_data, 1):
            topic_name = self._sanitize_filename(content["topic"])
            
            # Create topic folder
            if self.use_topic_folders:
                topic_folder = os.path.join(self.output_dir, f"Lecture_{idx}_{topic_name}")
                os.makedirs(topic_folder, exist_ok=True)
                base_path = os.path.join(topic_folder, topic_name)  # Full path including folder
            else:
                base_path = os.path.join(self.output_dir, f"{safe_subject}_{idx}_{topic_name}")
            
            files = {
                "topic": content["topic"],
                "unit": content["unit"],
                "folder": topic_folder if self.use_topic_folders else None,
                "files": {}
            }
            
            print(f"\n[Generating files for: {content['topic']}]")
            if self.use_topic_folders:
                print(f"   Folder: {os.path.basename(topic_folder)}/")
            
            # Generate PPT
            print("   [PPT] Creating PowerPoint...")
            ppt_file = self.generate_ppt(content, base_path)
            files["files"]["ppt"] = ppt_file
            
            # Generate PDF
            print("   [PDF] Creating PDF notes...")
            pdf_file = self.generate_pdf(content, base_path)
            files["files"]["pdf"] = pdf_file
            
            # Generate Audio with dynamics
            print("   [MP3] Creating audio lecture with voice dynamics...")
            audio_file = self.generate_audio_with_dynamics(content, base_path)
            files["files"]["audio"] = audio_file
            
            # Generate Video
            print("   [MP4] Creating video...")
            video_file = self.generate_video(content, base_path)
            files["files"]["video"] = video_file
            
            generated_files["files"].append(files)
        
        # Generate summary document
        print("\n[SUMMARY] Creating summary document...")
        summary_file = self.generate_summary(content_data, safe_subject, timestamp)
        generated_files["summary"] = summary_file
        
        return generated_files
    
    def generate_ppt(self, content, base_name):
        """Generate PowerPoint presentation"""
        filename = f"{base_name}.pptx"
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide layout
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = content["topic"]
        subtitle.text = f"{content['unit']}\nDifficulty: {content['difficulty']}"
        
        # Add content slides
        for slide_data in content["ppt_slides"]:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = slide_data["title"]
            
            tf = body.text_frame
            for bullet in slide_data["bullets"]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
        
        # Save
        prs.save(filename)
        return filename
    
    def generate_pdf(self, content, base_name):
        """Generate PDF notes"""
        filename = f"{base_name}.pdf"
        doc = SimpleDocTemplate(filename, pagesize=letter,
                              rightMargin=72, leftMargin=72,
                              topMargin=72, bottomMargin=18)
        
        story = []
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor='#667eea',
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor='#764ba2',
            spaceAfter=12,
            spaceBefore=12
        )
        
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['BodyText'],
            fontSize=11,
            alignment=TA_JUSTIFY,
            spaceAfter=12
        )
        
        # Title
        story.append(Paragraph(content["topic"], title_style))
        story.append(Spacer(1, 12))
        
        # Unit and difficulty
        story.append(Paragraph(f"<b>Unit:</b> {content['unit']}", body_style))
        story.append(Paragraph(f"<b>Difficulty:</b> {content['difficulty']}", body_style))
        story.append(Spacer(1, 20))
        
        # Learning objectives
        story.append(Paragraph("Learning Objectives", heading_style))
        for obj in content["learning_objectives"]:
            story.append(Paragraph(f"• {obj}", body_style))
        story.append(Spacer(1, 20))
        
        # PDF notes content
        notes_lines = content["pdf_notes"].split('\n')
        for line in notes_lines:
            if line.strip():
                if line.startswith('# '):
                    story.append(Paragraph(line[2:], title_style))
                elif line.startswith('## '):
                    story.append(Paragraph(line[3:], heading_style))
                elif line.startswith('**') and line.endswith('**'):
                    story.append(Paragraph(f"<b>{line[2:-2]}</b>", body_style))
                else:
                    story.append(Paragraph(line, body_style))
        
        doc.build(story)
        return filename

    
    def generate_audio_with_dynamics(self, content, base_name):
        """Generate audio with human-like voice using Microsoft Edge TTS (neural voices)"""
        filename = f"{base_name}.mp3"
        
        # Get the script - use clean text without SSML for edge-tts
        script = content["audio_script"]
        
        try:
            import asyncio
            import edge_tts
            
            print("      Using Microsoft neural voice (sounds more human)...")
            
            # Use Microsoft's neural voice (sounds very natural)
            # Options: en-US-AriaNeural (female), en-US-GuyNeural (male), en-US-JennyNeural (female)
            voice = "en-US-GuyNeural"  # Male voice, very natural
            
            # Generate audio asynchronously - edge-tts handles prosody automatically
            async def generate():
                communicate = edge_tts.Communicate(script, voice, rate="-5%", pitch="+0Hz")
                await communicate.save(filename)
            
            # Run the async function
            asyncio.run(generate())
            
            print(f"      ✅ Natural-sounding audio created")
            return filename
            
        except ImportError:
            print("      ⚠️ edge-tts not available, trying pyttsx3...")
            return self._generate_with_pyttsx3(content, base_name, script)
        except Exception as e:
            print(f"      ⚠️ Neural voice failed ({str(e)[:50]}), using fallback...")
            return self._generate_with_pyttsx3(content, base_name, script)
    
    def _generate_with_pyttsx3(self, content, base_name, script):
        """Fallback: Generate with pyttsx3 (offline but less natural)"""
        filename = f"{base_name}.mp3"
        
        try:
            import pyttsx3
            
            print("      Using pyttsx3 (offline voice)...")
            
            engine = pyttsx3.init()
            
            # Set properties for more natural speech
            voices = engine.getProperty('voices')
            # Try to use a better voice if available
            for voice in voices:
                if 'david' in voice.name.lower() or 'zira' in voice.name.lower():
                    engine.setProperty('voice', voice.id)
                    break
            
            # Adjust rate and volume for more natural sound
            engine.setProperty('rate', 160)  # Slightly slower than default
            engine.setProperty('volume', 0.9)
            
            # Save to file
            engine.save_to_file(script, filename)
            engine.runAndWait()
            
            print(f"      ✅ Audio created with offline voice")
            return filename
            
        except Exception as e:
            print(f"      ⚠️ pyttsx3 failed, using gTTS...")
            return self.generate_audio(content, base_name)
    
    def generate_audio(self, content, base_name):
        """Generate standard audio lecture (fallback)"""
        filename = f"{base_name}.mp3"
        
        # Use the actual script content
        script = content["audio_script"]
        
        # Generate audio
        try:
            tts = gTTS(text=script, lang='en', slow=False)
            tts.save(filename)
        except Exception as e:
            print(f"      ⚠️ Audio generation failed: {e}")
            # Create text file as fallback
            with open(filename.replace('.mp3', '_SCRIPT.txt'), 'w') as f:
                f.write(script)
            return filename.replace('.mp3', '_SCRIPT.txt')
        
        return filename
    
    def generate_video(self, content, base_name):
        """Generate actual MP4 video file - SIMPLE WORKING VERSION"""
        video_filename = f"{base_name}.mp4"
        
        # Create slide images
        frames_dir = os.path.join(os.path.dirname(base_name), "slides")
        os.makedirs(frames_dir, exist_ok=True)
        
        print("      Creating slide images...")
        slide_images = []
        for idx, slide in enumerate(content["ppt_slides"]):
            img_path = os.path.join(frames_dir, f"slide_{idx:02d}.png")
            self._create_slide_image(slide, img_path)
            slide_images.append(img_path)
        
        # Get audio file
        audio_file = f"{base_name}.mp3"
        
        if not os.path.exists(audio_file):
            print("      Audio file not found")
            return self._create_video_package(content, base_name, slide_images, audio_file)
        
        try:
            import cv2
            import numpy as np
            from pydub import AudioSegment
            import wave
            
            print("      Creating video with OpenCV...")
            
            # Get audio duration
            try:
                audio = AudioSegment.from_mp3(audio_file)
                total_duration_ms = len(audio)
                total_duration_sec = total_duration_ms / 1000.0
            except:
                # Fallback: 3 seconds per slide
                total_duration_sec = len(slide_images) * 3
            
            # Calculate FPS and frames per slide
            fps = 1  # 1 frame per second (since slides are static)
            duration_per_slide = total_duration_sec / len(slide_images)
            
            # Read first image to get dimensions
            first_img = cv2.imread(slide_images[0])
            height, width, layers = first_img.shape
            
            # Create video writer (without audio first)
            temp_video = video_filename.replace('.mp4', '_temp.mp4')
            fourcc = cv2.VideoWriter_fourcc(*'mp4v')
            video_writer = cv2.VideoWriter(temp_video, fourcc, fps, (width, height))
            
            # Write each slide for its duration
            for img_path in slide_images:
                img = cv2.imread(img_path)
                frames_to_write = int(duration_per_slide * fps)
                for _ in range(max(1, frames_to_write)):
                    video_writer.write(img)
            
            video_writer.release()
            
            # Now combine video with audio using ffmpeg
            print("      Adding audio to video...")
            import subprocess
            
            cmd = [
                'ffmpeg', '-y',
                '-i', temp_video,
                '-i', audio_file,
                '-c:v', 'copy',
                '-c:a', 'aac',
                '-shortest',
                video_filename
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            # Clean up temp file
            if os.path.exists(temp_video):
                os.remove(temp_video)
            
            if result.returncode == 0 and os.path.exists(video_filename):
                file_size = os.path.getsize(video_filename) / (1024 * 1024)  # MB
                print(f"      ✅ Video created: {os.path.basename(video_filename)} ({file_size:.1f} MB)")
                return video_filename
            else:
                raise Exception("ffmpeg not available")
                
        except ImportError as e:
            print(f"      ⚠️ Missing library: {str(e)[:50]}")
            print("         Creating video package instead")
            return self._create_video_package(content, base_name, slide_images, audio_file)
        except FileNotFoundError:
            print("      ⚠️ ffmpeg not found")
            print("         Creating video package instead")
            return self._create_video_package(content, base_name, slide_images, audio_file)
        except Exception as e:
            print(f"      ⚠️ Video creation failed: {str(e)[:80]}")
            print("         Creating video package instead")
            return self._create_video_package(content, base_name, slide_images, audio_file)
    
    def _create_video_package(self, content, base_name, slide_images, audio_file):
        """Create a package with slides and audio for manual video creation"""
        package_dir = os.path.join(self.output_dir, f"{base_name}_VIDEO_PACKAGE")
        os.makedirs(package_dir, exist_ok=True)
        
        # Copy audio
        import shutil
        if os.path.exists(audio_file):
            shutil.copy(audio_file, os.path.join(package_dir, "audio.mp3"))
        
        # Copy slides
        slides_dir = os.path.join(package_dir, "slides")
        os.makedirs(slides_dir, exist_ok=True)
        for idx, img in enumerate(slide_images):
            if os.path.exists(img):
                shutil.copy(img, os.path.join(slides_dir, f"slide_{idx:02d}.png"))
        
        # Create instructions
        instructions = os.path.join(package_dir, "HOW_TO_CREATE_VIDEO.txt")
        with open(instructions, 'w') as f:
            f.write(f"VIDEO PACKAGE: {content['topic']}\n\n")
            f.write("This package contains everything needed to create a video:\n\n")
            f.write("FILES:\n")
            f.write("- audio.mp3: The lecture audio\n")
            f.write(f"- slides/: {len(slide_images)} slide images\n\n")
            f.write("OPTION 1: Use ffmpeg (command line)\n")
            f.write("Install ffmpeg, then run:\n")
            f.write("  ffmpeg -framerate 1/3 -i slides/slide_%02d.png -i audio.mp3 ")
            f.write("-c:v libx264 -c:a aac -shortest -pix_fmt yuv420p output.mp4\n\n")
            f.write("OPTION 2: Use video editing software\n")
            f.write("- Import all slides from slides/ folder\n")
            f.write("- Import audio.mp3\n")
            f.write("- Arrange slides to match audio duration\n")
            f.write("- Export as MP4\n\n")
            f.write("Recommended software:\n")
            f.write("- Windows: Windows Movie Maker, DaVinci Resolve\n")
            f.write("- Mac: iMovie, Final Cut Pro\n")
            f.write("- Online: Canva, Kapwing\n")
        
        print(f"      ✅ Video package created: {os.path.basename(package_dir)}/")
        print(f"         Contains {len(slide_images)} slides + audio")
        print(f"         See HOW_TO_CREATE_VIDEO.txt for instructions")
        
        return package_dir
    
    def _create_video_plan(self, content, base_name, video_filename, slide_images=None):
        """Create video plan as fallback"""
        if slide_images is None:
            # Create slide images
            frames_dir = os.path.join(self.output_dir, "temp_frames")
            os.makedirs(frames_dir, exist_ok=True)
            
            slide_images = []
            for idx, slide in enumerate(content["ppt_slides"]):
                img_path = os.path.join(frames_dir, f"slide_{idx:02d}.png")
                self._create_slide_image(slide, img_path)
                slide_images.append(img_path)
        
        video_plan_file = video_filename.replace('.mp4', '_VIDEO_PLAN.txt')
        with open(video_plan_file, 'w') as f:
            f.write(f"VIDEO CONTENT PLAN: {content['topic']}\n\n")
            if 'video_content' in content:
                f.write(f"Duration: {content['video_content']['duration']}\n\n")
                f.write("SEGMENTS:\n\n")
                for segment in content['video_content']['segments']:
                    f.write(f"{segment['time']} - {segment['title']}\n")
                    f.write(f"Content: {segment['content']}\n")
                    f.write(f"Visuals: {segment['visuals']}\n\n")
            f.write(f"\nSlide images created: {len(slide_images)} slides\n")
            f.write(f"Location: {os.path.dirname(slide_images[0]) if slide_images else 'N/A'}\n")
            f.write(f"\nTo create MP4 video:\n")
            f.write(f"1. Install ffmpeg: https://ffmpeg.org/download.html\n")
            f.write(f"2. Or use video editing software (iMovie, DaVinci Resolve, etc.)\n")
        
        return video_plan_file
    
    def _create_slide_image(self, slide_data, output_path):
        """Create beautiful, aesthetic slide image with modern design"""
        try:
            from PIL import ImageDraw, ImageFont, ImageFilter
            
            # Create image with gradient background
            img = Image.new('RGB', (1920, 1080), color='white')
            draw = ImageDraw.Draw(img)
            
            # Create gradient background
            for y in range(1080):
                # Gradient from purple to blue
                r = int(102 + (66 - 102) * y / 1080)
                g = int(126 + (135 - 126) * y / 1080)
                b = int(234 + (245 - 234) * y / 1080)
                draw.line([(0, y), (1920, y)], fill=(r, g, b))
            
            # Add semi-transparent overlay for better text readability
            overlay = Image.new('RGBA', (1920, 1080), (255, 255, 255, 30))
            img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
            
            draw = ImageDraw.Draw(img)
            
            # Try to load nice fonts
            try:
                title_font = ImageFont.truetype("arial.ttf", 80)
                text_font = ImageFont.truetype("arial.ttf", 45)
            except:
                try:
                    title_font = ImageFont.truetype("C:\\Windows\\Fonts\\arial.ttf", 80)
                    text_font = ImageFont.truetype("C:\\Windows\\Fonts\\arial.ttf", 45)
                except:
                    title_font = ImageFont.load_default()
                    text_font = ImageFont.load_default()
            
            # Draw decorative top bar
            draw.rectangle([(0, 0), (1920, 20)], fill=(255, 255, 255, 200))
            
            # Draw title with shadow for depth
            title = slide_data['title']
            # Shadow
            draw.text((152, 152), title, fill=(0, 0, 0, 100), font=title_font)
            # Main title
            draw.text((150, 150), title, fill='white', font=title_font)
            
            # Draw decorative line under title
            draw.rectangle([(150, 260), (1770, 265)], fill='white')
            
            # Draw bullets with icons
            y_pos = 350
            for i, bullet in enumerate(slide_data['bullets'][:6]):  # Max 6 bullets
                # Bullet point circle
                draw.ellipse([(180, y_pos + 10), (210, y_pos + 40)], fill='white')
                
                # Bullet text with shadow
                bullet_text = bullet[:100] if len(bullet) > 100 else bullet  # Limit length
                # Shadow
                draw.text((252, y_pos + 2), bullet_text, fill=(0, 0, 0, 80), font=text_font)
                # Main text
                draw.text((250, y_pos), bullet_text, fill='white', font=text_font)
                
                y_pos += 100
            
            # Add footer with slide number
            footer_text = f"Slide {slide_data['slide_number']}"
            draw.text((1700, 1020), footer_text, fill='white', font=text_font)
            
            # Save
            img.save(output_path, quality=95)
            
        except Exception as e:
            print(f"      ⚠️ Aesthetic slide creation failed: {e}")
            # Fallback to simple slide
            self._create_simple_slide(slide_data, output_path)
    
    def _create_simple_slide(self, slide_data, output_path):
        """Fallback: Create simple slide"""
        try:
            img = Image.new('RGB', (1920, 1080), color='#667eea')
            draw = ImageDraw.Draw(img)
            
            try:
                title_font = ImageFont.truetype("arial.ttf", 60)
                text_font = ImageFont.truetype("arial.ttf", 40)
            except:
                title_font = ImageFont.load_default()
                text_font = ImageFont.load_default()
            
            # Draw title
            draw.text((100, 100), slide_data['title'], fill='white', font=title_font)
            
            # Draw bullets
            y_pos = 300
            for bullet in slide_data['bullets']:
                bullet_text = f"• {bullet[:80]}"
                draw.text((150, y_pos), bullet_text, fill='white', font=text_font)
                y_pos += 100
            
            img.save(output_path)
            
        except Exception as e:
            print(f"      ⚠️ Simple slide creation failed: {e}")

    
    def generate_summary(self, content_data, subject_name, timestamp):
        """Generate summary PDF with all topics"""
        filename = os.path.join(self.output_dir, f"{subject_name}_SUMMARY_{timestamp}.pdf")
        doc = SimpleDocTemplate(filename, pagesize=letter,
                              rightMargin=72, leftMargin=72,
                              topMargin=72, bottomMargin=18)
        
        story = []
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'Title',
            parent=styles['Heading1'],
            fontSize=24,
            textColor='#667eea',
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'Heading',
            parent=styles['Heading2'],
            fontSize=14,
            textColor='#764ba2',
            spaceAfter=12
        )
        
        body_style = styles['BodyText']
        
        # Title
        story.append(Paragraph(f"Course Content Summary: {subject_name}", title_style))
        story.append(Paragraph(f"Generated: {timestamp}", body_style))
        story.append(Spacer(1, 30))
        
        # Table of contents
        story.append(Paragraph("Table of Contents", heading_style))
        for idx, content in enumerate(content_data, 1):
            story.append(Paragraph(f"{idx}. {content['topic']} ({content['unit']})", body_style))
        
        story.append(PageBreak())
        
        # Content summaries
        for idx, content in enumerate(content_data, 1):
            story.append(Paragraph(f"{idx}. {content['topic']}", heading_style))
            story.append(Paragraph(f"<b>Unit:</b> {content['unit']}", body_style))
            story.append(Paragraph(f"<b>Difficulty:</b> {content['difficulty']}", body_style))
            story.append(Spacer(1, 12))
            
            story.append(Paragraph("<b>Learning Objectives:</b>", body_style))
            for obj in content["learning_objectives"]:
                story.append(Paragraph(f"• {obj}", body_style))
            
            story.append(Spacer(1, 12))
            story.append(Paragraph(f"<b>Slides:</b> {len(content['ppt_slides'])} slides", body_style))
            story.append(Spacer(1, 20))
            
            if idx < len(content_data):
                story.append(PageBreak())
        
        doc.build(story)
        return filename
    
    def _sanitize_filename(self, name):
        """Clean filename for safe file system usage"""
        # Remove special characters
        safe = "".join(c for c in name if c.isalnum() or c in (' ', '-', '_'))
        safe = safe.replace(' ', '_')
        return safe[:50]  # Limit length


def main():
    """Test file generation"""
    from course_content_generator import CourseContentGenerator, ContentInput
    
    # Sample input
    input_data = ContentInput(
        syllabus_text="""
        Unit 1: Python Basics
        - Variables and Data Types
        - Control Flow
        
        Unit 2: Functions
        - Function Definition
        - Parameters and Arguments
        """,
        course_outline_text="Introduction to Python programming",
        timetable_text="Mon/Wed/Fri 10:00 AM",
        subject_name="Python Programming",
        course_duration="10 weeks",
        class_duration=60,
        mode="Weekly"
    )
    
    # Generate content
    print("🚀 Generating course content...")
    generator = CourseContentGenerator()
    result = generator.generate(input_data)
    
    # Generate files
    print("\n📁 Creating files...")
    file_gen = FileGenerator()
    files = file_gen.generate_all(result["content"], result["subject"])
    
    print("\n✅ Files generated successfully!")
    print(f"\n📂 Output directory: {file_gen.output_dir}")
    print("\n📄 Generated files:")
    
    for item in files["files"]:
        print(f"\n  Topic: {item['topic']}")
        for file_type, filepath in item["files"].items():
            print(f"    {file_type.upper()}: {os.path.basename(filepath)}")
    
    print(f"\n  Summary: {os.path.basename(files['summary'])}")
    print("\n" + "="*60)


if __name__ == "__main__":
    main()
